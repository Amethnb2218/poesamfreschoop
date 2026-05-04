#!/usr/bin/env python3
"""
FreshCoop - Generateur de dossier de candidature POESAM 2026
Genere : PDF, Word (DOCX), PowerPoint (PPTX)
"""
import os
from fpdf import FPDF
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DIR = os.path.dirname(__file__)

# ============================================================
# CONTENT DATA
# ============================================================

GREEN = (34, 139, 34)
DARK = (30, 30, 30)
GRAY = (100, 100, 100)
WHITE = (255, 255, 255)
LIGHT_GREEN = (230, 245, 230)

RESUME_EXECUTIF = (
    "FreshCoop est une plateforme agri-tech cooperative concue au Senegal pour resoudre "
    "trois defis majeurs auxquels font face les petits producteurs agricoles : les pertes "
    "post-recolte massives (estimees a 40%), l'absence de tracabilite des transactions, "
    "et un acces au credit formel quasi inexistant (moins de 5%)."
)

RESUME_PILIERS = [
    "Des micro-hubs frigorifiques solaires partages pour reduire les pertes post-recolte",
    "Un systeme d'intelligence marche qui oriente chaque lot vers le meilleur debouche",
    "Une preuve economique portable permettant aux producteurs d'acceder au credit formel",
]

RESUME_SUITE = (
    "FreshCoop connecte producteurs, acheteurs B2B, transporteurs, agents de terrain et "
    "partenaires financiers dans un ecosysteme transparent, inclusif et consent-driven. "
    "La plateforme est disponible en version web, mobile (iOS/Android) et USSD (*384*FRES#) "
    "pour garantir l'accessibilite a tous, y compris les 70% de producteurs sans smartphone."
)

RESUME_CANDIDATURE = (
    "Candidate au programme POESAM 2026, FreshCoop vise a demontrer un modele economique "
    "viable, scalable et a fort impact social dans les regions de Dakar, Thies, Saint-Louis "
    "et Casamance."
)

PROBLEME_INTRO = (
    "L'agriculture senegalaise emploie plus de 60% de la population active, mais les petits "
    "producteurs restent pieges dans un cercle vicieux de pertes, d'invisibilite economique "
    "et d'exclusion financiere."
)

PROBLEMES = [
    ("Pertes post-recolte massives",
     "Au Senegal, pres de 40% des fruits et legumes produits sont perdus apres la recolte. "
     "L'absence d'infrastructures de stockage frigorifique, la logistique defaillante et le "
     "manque d'acces rapide aux canaux de vente entrainent la deterioration des produits "
     "avant qu'ils n'atteignent le consommateur."),
    ("Absence de tracabilite",
     "Les transactions dans les marches agricoles informels ne laissent aucune trace "
     "exploitable. Les producteurs n'ont aucun historique prouvable de leur activite "
     "commerciale, ce qui les rend invisibles aux yeux des institutions financieres."),
    ("Exclusion financiere",
     "Moins de 5% des petits producteurs senegalais ont acces au credit formel. Sans "
     "preuves d'activite economique, les banques et les systemes financiers decentralises "
     "(SFD) ne peuvent evaluer leur solvabilite."),
]

SOLUTION_INTRO = (
    "FreshCoop repond a ces trois defis avec une plateforme integree qui couvre "
    "l'ensemble de la chaine de valeur agricole, de la recolte au financement."
)

SOLUTIONS = [
    ("Micro-hubs frigorifiques solaires",
     "Des unites de stockage frigorifique alimentees a l'energie solaire, partagees entre "
     "plusieurs producteurs d'une meme zone. Chaque hub est equipe de capteurs IoT "
     "surveillant en temps reel la temperature, la capacite, le niveau de batterie et "
     "l'etat des stocks."),
    ("Intelligence marche",
     "Un algorithme d'orientation intelligent analyse pour chaque lot : le prix du marche, "
     "la demande locale, le score de confiance des acheteurs et la fenetre de consommation "
     "optimale. Il recommande le meilleur debouche pour maximiser les revenus du producteur."),
    ("Preuve economique portable",
     "Chaque transaction realisee via FreshCoop genere une preuve economique numerique "
     "verifiable. Les producteurs constituent un dossier de bancabilite portable : "
     "historique de ventes, relations acheteurs, paiements recus, volumes traites."),
    ("Acces inclusif via USSD",
     "Pour les 70% de producteurs sans smartphone, FreshCoop propose un acces USSD "
     "(*384*FRES#) avec un menu en Wolof et Pulaar : consultation des prix du jour, "
     "mise a jour des stocks, enregistrement de ventes et reception de paiements."),
]

MARCHE_TAILLE = (
    "Le marche agricole senegalais represente plus de 3 000 milliards FCFA par an. "
    "Le segment des fruits et legumes, notre cible primaire, pese environ 800 milliards FCFA. "
    "Avec des pertes post-recolte de 40%, la valeur detruite chaque annee depasse "
    "300 milliards FCFA."
)

CONCURRENTS = [
    ("Mlouma / Afrimalin", "Places de marche en ligne sans tracabilite ni stockage froid"),
    ("ColdHubs (Nigeria)", "Stockage solaire mais sans marketplace ni preuve economique"),
    ("Twiga Foods (Kenya)", "B2B supply chain mais sans dimension d'inclusion financiere"),
    ("FarmCrowdy", "Financement participatif sans intelligence marche"),
]

AVANTAGES = [
    "Pas de wallet proprietaire : paiements via Orange Money, Wave, PayDunya",
    "Donnees sous controle du producteur avec consentement explicite et revocable",
    "Tracabilite complete lot par lot avec QR code, photos et capteurs IoT",
    "Scoring de credit explicable base sur des transactions reelles",
    "Accessible sans smartphone via USSD en langues locales",
]

ROLES = [
    ("Producteur (Agriculteur)", "Publication de produits, creation de lots avec QR code, suivi des commandes, intelligence marche, export de preuve economique"),
    ("Acheteur B2B", "Reservation de lots en gros, tracabilite champ-a-marche, rachat facilite, specifications qualite"),
    ("Client (Marketplace)", "Navigation par categorie, panier d'achat, messagerie vendeur, suivi de commande"),
    ("Transporteur", "Gestion des itineraires, monitoring des hubs, alertes de risque"),
    ("Agent de terrain", "Confirmation de commandes hors-ligne, coordination avec producteurs et transporteurs"),
    ("Partenaire financier", "Consultation des preuves economiques consenties, scoring de bancabilite"),
    ("Administrateur", "Gestion des utilisateurs, validation des dossiers, KPIs plateforme en temps reel"),
]

MODULES = [
    ("Alertes anti-gaspillage", "Surveillance DLC en temps reel, alertes automatiques pour vente rapide des lots a risque"),
    ("Paiements orchestres", "Orange Money, Wave, PayDunya, banques. Codes recus verifiables"),
    ("Dossiers et attestations", "Soumission de documents, validation et delivrance d'attestations"),
    ("Dashboard d'impact", "KPIs POESAM en temps reel : pertes evitees, revenus additionnels, CO2 economise"),
]

EQUIPE = [
    ("Direction Generale / Fondateur", "Vision strategique, connaissance du terrain agricole senegalais, relations institutionnelles POESAM"),
    ("CTO / Lead Developpement", "Architecture full-stack (React, Node.js, React Native/Expo), integration IoT, securite des donnees"),
    ("Responsable Operations Terrain", "Coordination avec les cooperatives, deploiement des hubs, formation des agents terrain"),
    ("Responsable Partenariats Financiers", "Relations avec les SFD, banques et fintechs, conception du module de preuve economique"),
    ("Data Analyst / Impact", "Modelisation des KPIs POESAM, mesure d'impact, reporting aux bailleurs"),
    ("UX / Design", "Conception d'interfaces accessibles, adaptation multilingue (Wolof, Pulaar, Francais)"),
]

MODELE_REVENUS = [
    "Commission sur transactions : 2-5% par commande confirmee et livree",
    "Abonnement hub : 5 000 - 15 000 FCFA/mois pour le stockage frigorifique",
    "Services premium B2B : visibilite renforcee, analytics avances, acces prioritaire",
    "Frais de certification : generation et export de preuves economiques",
    "Data insights : analyses de marche anonymisees pour ministeres et ONG (a terme)",
]

PREVISIONS = [
    ("Revenus commissions", "8", "45", "180"),
    ("Abonnements hubs", "5", "25", "90"),
    ("Services B2B premium", "2", "15", "60"),
    ("Certification/preuves", "1", "8", "35"),
    ("TOTAL REVENUS", "16", "93", "365"),
    ("Couts infrastructure", "-12", "-30", "-80"),
    ("Salaires equipe", "-18", "-35", "-60"),
    ("Marketing / acquisition", "-5", "-15", "-30"),
    ("Operations terrain", "-8", "-20", "-40"),
    ("Technologie / hebergement", "-4", "-8", "-15"),
    ("TOTAL CHARGES", "-47", "-108", "-225"),
    ("RESULTAT NET", "-31", "-15", "+140"),
]

TRESORERIE = [
    ("Tresorerie debut", "0", "19", "54"),
    ("Financement POESAM", "50", "0", "0"),
    ("Autres financements", "0", "50", "0"),
    ("Resultat d'exploitation", "-31", "-15", "+140"),
    ("Tresorerie fin", "19", "54", "194"),
]

FINANCEMENT_REPARTITION = [
    ("Deploiement hubs solaires (x4)", "20 M FCFA", "40%"),
    ("Developpement plateforme", "10 M FCFA", "20%"),
    ("Operations terrain et formation", "8 M FCFA", "16%"),
    ("Marketing et acquisition", "5 M FCFA", "10%"),
    ("Fonds de roulement", "4 M FCFA", "8%"),
    ("Conformite et legal", "3 M FCFA", "6%"),
]

IMPACT_SOCIAL = [
    ("Producteurs connectes", "100", "An 1"),
    ("Dont femmes productrices", "> 40%", "An 1"),
    ("Cooperatives actives", "3", "An 1"),
    ("Acces credit formel", "30 producteurs", "An 1"),
    ("Revenu additionnel moyen", "+25%", "An 2"),
    ("Emplois directs crees", "15", "An 1"),
    ("Emplois indirects", "50+", "An 2"),
]

IMPACT_ENVIRONNEMENTAL = [
    ("Reduction pertes post-recolte", "-30%", "An 1"),
    ("Tonnes de CO2 evitees", "50 tonnes", "An 1"),
    ("Energie solaire produite", "100% renouvelable", "Continu"),
    ("Kg de produits traces", "150 000 kg", "An 1"),
    ("Alertes anti-gaspillage", "500+", "An 1"),
]

ODD = [
    ("ODD 1", "Pas de pauvrete", "Acces au credit et amelioration des revenus"),
    ("ODD 2", "Faim zero", "Reduction des pertes alimentaires"),
    ("ODD 5", "Egalite des sexes", "Inclusion des femmes productrices (>40%)"),
    ("ODD 7", "Energie propre", "Hubs solaires 100% renouvelable"),
    ("ODD 12", "Consommation responsable", "Tracabilite et anti-gaspillage"),
    ("ODD 13", "Action climatique", "Reduction de l'empreinte carbone"),
]

INNOVATION = [
    "IA predictive pour l'optimisation des prix et la prevision de la demande",
    "IoT avance avec capteurs de qualite (humidite, ethylene) dans les hubs",
    "Blockchain pour la certification d'origine des produits",
    "Machine learning pour le scoring de credit plus precis",
]

CROISSANCE = [
    ("Phase 1 (An 1)", "Pilote dans 4 regions avec 3 cooperatives partenaires"),
    ("Phase 2 (An 2)", "Extension a 6 regions, partenariats avec SFD et banques"),
    ("Phase 3 (An 3)", "Couverture nationale, ouverture a l'export sous-regional"),
    ("Phase 4 (An 4+)", "Expansion Afrique de l'Ouest (Mali, Guinee, Gambie)"),
]

PARTENARIATS = [
    ("Operateurs telecom", "Orange, Free pour l'infrastructure USSD"),
    ("Fintechs", "Wave, PayDunya pour l'orchestration des paiements"),
    ("Institutions", "POESAM, ministere de l'Agriculture, SAED"),
    ("ONG", "FAO, PAM pour le co-financement et l'expertise terrain"),
    ("Banques/SFD", "CMS, PAMECAS pour le credit aux producteurs"),
]

OBJECTIFS_PILOTE = [
    ("Producteurs inscrits", "100", "12 mois"),
    ("Lots traces avec QR", "300", "12 mois"),
    ("Commandes B2B confirmees", "150", "12 mois"),
    ("Paiements via partenaires", "100", "12 mois"),
    ("Preuves economiques exportees", "100", "12 mois"),
    ("Acheteurs B2B actifs", "20", "12 mois"),
    ("Cooperatives pilotes", "3", "12 mois"),
]

TOC = [
    ("1", "Resume executif"),
    ("2", "Probleme identifie"),
    ("3", "Solution proposee : FreshCoop"),
    ("4", "Analyse du marche et des concurrents"),
    ("5", "Produit et fonctionnalites"),
    ("6", "Equipe et competences cles"),
    ("7", "Modele economique"),
    ("8", "Previsions financieres sur 3 ans"),
    ("9", "Utilisation du financement"),
    ("10", "Impact social et environnemental"),
    ("11", "Strategie d'innovation et de croissance"),
    ("12", "Objectifs du pilote POESAM"),
    ("13", "Conclusion"),
]


# ============================================================
# PDF GENERATION
# ============================================================

class FreshCoopPDF(FPDF):
    def header(self):
        if self.page_no() > 1:
            self.set_font("Helvetica", "B", 9)
            self.set_text_color(*GREEN)
            self.cell(0, 8, "FreshCoop - Dossier de Candidature POESAM 2026", align="L")
            self.ln(4)
            self.set_draw_color(*GREEN)
            self.set_line_width(0.5)
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*GRAY)
        if self.page_no() > 1:
            self.cell(0, 10, f"FreshCoop  |  Confidentiel  |  Page {self.page_no()}", align="C")

    def cover_page(self):
        self.add_page()
        self.ln(40)
        self.set_fill_color(*GREEN)
        self.rect(0, 30, 210, 80, "F")
        self.set_y(45)
        self.set_font("Helvetica", "B", 36)
        self.set_text_color(*WHITE)
        self.cell(0, 18, "FRESHCOOP", align="C", new_x="LMARGIN", new_y="NEXT")
        self.set_font("Helvetica", "", 14)
        self.cell(0, 10, "La plateforme cooperative qui connecte", align="C", new_x="LMARGIN", new_y="NEXT")
        self.cell(0, 8, "les producteurs aux marches", align="C", new_x="LMARGIN", new_y="NEXT")
        self.ln(10)
        self.set_font("Helvetica", "B", 11)
        self.cell(0, 8, "Micro-hubs solaires  |  Intelligence marche  |  Preuve economique", align="C", new_x="LMARGIN", new_y="NEXT")

        self.set_y(130)
        self.set_text_color(*DARK)
        self.set_font("Helvetica", "B", 16)
        self.cell(0, 12, "Dossier de Candidature", align="C", new_x="LMARGIN", new_y="NEXT")
        self.set_font("Helvetica", "", 13)
        self.set_text_color(*GRAY)
        self.cell(0, 10, "Programme POESAM 2026", align="C", new_x="LMARGIN", new_y="NEXT")
        self.ln(15)

        info = [
            ("Porteur du projet", "FreshCoop SAS"),
            ("Localisation", "Dakar, Senegal"),
            ("Secteur", "AgriTech / FinTech / Impact Social"),
            ("Date", "Mai 2026"),
            ("Contact", "contact@freshcoop.sn"),
        ]
        self.set_font("Helvetica", "", 11)
        for label, value in info:
            self.set_text_color(*GRAY)
            self.cell(75, 8, label, align="R")
            self.set_text_color(*DARK)
            self.set_font("Helvetica", "B", 11)
            self.cell(5, 8, "")
            self.cell(0, 8, value, align="L", new_x="LMARGIN", new_y="NEXT")
            self.set_font("Helvetica", "", 11)

    def section_title(self, num, title):
        self.ln(6)
        self.set_fill_color(*GREEN)
        self.set_text_color(*WHITE)
        self.set_font("Helvetica", "B", 14)
        self.cell(0, 12, f"  {num}. {title}", fill=True, new_x="LMARGIN", new_y="NEXT")
        self.ln(4)
        self.set_text_color(*DARK)

    def sub_title(self, title):
        self.ln(3)
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(*GREEN)
        self.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
        self.set_text_color(*DARK)
        self.set_font("Helvetica", "", 10)

    def body_text(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*DARK)
        self.multi_cell(0, 6, text)
        self.ln(2)

    def bullet(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(*DARK)
        self.multi_cell(0, 6, f"  > {text}", new_x="LMARGIN", new_y="NEXT")

    def kpi_table(self, rows):
        self.ln(2)
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(*GREEN)
        self.set_text_color(*WHITE)
        self.cell(80, 8, "  Indicateur", fill=True)
        self.cell(50, 8, "Objectif", fill=True, align="C")
        self.cell(50, 8, "Horizon", fill=True, align="C", new_x="LMARGIN", new_y="NEXT")
        self.set_text_color(*DARK)
        self.set_font("Helvetica", "", 9)
        fill = False
        for ind, obj, hor in rows:
            if fill:
                self.set_fill_color(*LIGHT_GREEN)
            self.cell(80, 7, f"  {ind}", fill=fill)
            self.cell(50, 7, obj, fill=fill, align="C")
            self.cell(50, 7, hor, fill=fill, align="C", new_x="LMARGIN", new_y="NEXT")
            fill = not fill
        self.ln(3)

    def finance_table(self, title, rows, cols=None):
        if cols is None:
            cols = ["Poste", "Annee 1", "Annee 2", "Annee 3"]
        self.sub_title(title)
        self.ln(1)
        ncols = len(cols)
        w_first = 60
        w_other = (180 - w_first) // (ncols - 1) if ncols > 1 else 0
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(*GREEN)
        self.set_text_color(*WHITE)
        for i, c in enumerate(cols):
            w = w_first if i == 0 else w_other
            al = "L" if i == 0 else "C"
            kw = {}
            if i == ncols - 1:
                kw = {"new_x": "LMARGIN", "new_y": "NEXT"}
            self.cell(w, 8, f"  {c}" if i == 0 else c, fill=True, align=al, **kw)
        self.set_text_color(*DARK)
        self.set_font("Helvetica", "", 9)
        fill = False
        for row in rows:
            bold = row[0].startswith("TOTAL") or row[0].startswith("RESULTAT") or row[0] == "Tresorerie fin"
            if bold:
                self.set_font("Helvetica", "B", 9)
            if fill:
                self.set_fill_color(*LIGHT_GREEN)
            for i, val in enumerate(row):
                w = w_first if i == 0 else w_other
                al = "L" if i == 0 else "C"
                kw = {}
                if i == len(row) - 1:
                    kw = {"new_x": "LMARGIN", "new_y": "NEXT"}
                self.cell(w, 7, f"  {val}" if i == 0 else val, fill=fill, align=al, **kw)
            if bold:
                self.set_font("Helvetica", "", 9)
            fill = not fill
        self.ln(4)


def build_pdf():
    pdf = FreshCoopPDF()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Cover
    pdf.cover_page()

    # TOC
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.set_text_color(*GREEN)
    pdf.cell(0, 14, "Table des matieres", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    for num, title in TOC:
        pdf.set_font("Helvetica", "B", 11)
        pdf.set_text_color(*DARK)
        pdf.cell(12, 8, num + ".")
        pdf.set_font("Helvetica", "", 11)
        pdf.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    # 1. Resume executif
    pdf.add_page()
    pdf.section_title("1", "Resume executif")
    pdf.body_text(RESUME_EXECUTIF)
    pdf.body_text("Notre solution repose sur trois piliers complementaires :")
    for p in RESUME_PILIERS:
        pdf.bullet(p)
    pdf.ln(2)
    pdf.body_text(RESUME_SUITE)
    pdf.body_text(RESUME_CANDIDATURE)

    # 2. Probleme
    pdf.section_title("2", "Probleme identifie")
    pdf.body_text(PROBLEME_INTRO)
    for i, (t, d) in enumerate(PROBLEMES, 1):
        pdf.sub_title(f"2.{i} {t}")
        pdf.body_text(d)

    # 3. Solution
    pdf.add_page()
    pdf.section_title("3", "Solution proposee : FreshCoop")
    pdf.body_text(SOLUTION_INTRO)
    for i, (t, d) in enumerate(SOLUTIONS, 1):
        pdf.sub_title(f"3.{i} {t}")
        pdf.body_text(d)

    # 4. Analyse du marche
    pdf.add_page()
    pdf.section_title("4", "Analyse du marche et des concurrents")
    pdf.sub_title("4.1 Taille du marche")
    pdf.body_text(MARCHE_TAILLE)
    pdf.sub_title("4.2 Paysage concurrentiel")
    for name, desc in CONCURRENTS:
        pdf.bullet(f"{name} : {desc}")
    pdf.ln(2)
    pdf.body_text(
        "FreshCoop est le seul acteur qui integre stockage froid solaire, intelligence marche, "
        "tracabilite complete et preuve economique portable dans une seule plateforme."
    )
    pdf.sub_title("4.3 Avantage concurrentiel")
    for a in AVANTAGES:
        pdf.bullet(a)

    # 5. Produit
    pdf.add_page()
    pdf.section_title("5", "Produit et fonctionnalites")
    pdf.sub_title("5.1 Plateforme web (React 19 + Vite + Node.js)")
    pdf.body_text("Application web complete avec tableau de bord adapte a chaque role.")
    pdf.sub_title("5.2 Application mobile (React Native / Expo)")
    pdf.body_text("Application native iOS et Android avec QR code, notifications push et GPS.")
    pdf.sub_title("5.3 Fonctionnalites par role")
    for role, desc in ROLES:
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*GREEN)
        pdf.cell(0, 7, role, new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(*DARK)
        pdf.set_font("Helvetica", "", 9)
        pdf.multi_cell(0, 5.5, desc)
        pdf.ln(2)
    pdf.sub_title("5.4 Modules transversaux")
    for mod, desc in MODULES:
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*GREEN)
        pdf.cell(0, 7, mod, new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(*DARK)
        pdf.set_font("Helvetica", "", 9)
        pdf.multi_cell(0, 5.5, desc)
        pdf.ln(2)

    # 6. Equipe
    pdf.add_page()
    pdf.section_title("6", "Equipe et competences cles")
    pdf.body_text("L'equipe FreshCoop reunit des competences complementaires :")
    for role, desc in EQUIPE:
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*GREEN)
        pdf.cell(0, 7, role, new_x="LMARGIN", new_y="NEXT")
        pdf.set_text_color(*DARK)
        pdf.set_font("Helvetica", "", 9)
        pdf.multi_cell(0, 5.5, desc)
        pdf.ln(2)

    # 7. Modele economique
    pdf.section_title("7", "Modele economique")
    pdf.body_text("FreshCoop adopte un modele de revenus diversifie :")
    for r in MODELE_REVENUS:
        pdf.bullet(r)
    pdf.ln(2)
    pdf.body_text(
        "Point cle : FreshCoop n'est PAS un portefeuille electronique. Les paiements sont "
        "orchestres par des partenaires agrees (Orange Money, Wave, PayDunya)."
    )

    # 8. Previsions financieres
    pdf.add_page()
    pdf.section_title("8", "Previsions financieres sur 3 ans")
    pdf.sub_title("8.1 Hypotheses cles")
    pdf.bullet("Annee 1 (Pilote) : 100 producteurs, 3 cooperatives, 4 regions")
    pdf.bullet("Annee 2 (Croissance) : 500 producteurs, 12 cooperatives, 6 regions")
    pdf.bullet("Annee 3 (Echelle) : 2 000 producteurs, 40 cooperatives, 10 regions")
    pdf.ln(2)
    pdf.finance_table("8.2 Compte de resultat previsionnel (M FCFA)", PREVISIONS)
    pdf.finance_table("8.3 Tresorerie previsionnelle (M FCFA)", TRESORERIE)
    pdf.body_text(
        "Le point mort (break-even) est atteint au cours de l'Annee 3, avec une marge "
        "nette positive de 38%."
    )

    # 9. Utilisation du financement
    pdf.add_page()
    pdf.section_title("9", "Utilisation prevue du financement")
    pdf.body_text("Le financement POESAM (50 millions FCFA) sera alloue comme suit :")
    pdf.finance_table("Repartition du financement", FINANCEMENT_REPARTITION,
                      cols=["Poste", "Montant", "Part"])

    # 10. Impact
    pdf.section_title("10", "Impact social et environnemental")
    pdf.sub_title("10.1 Impact social mesurable")
    pdf.kpi_table(IMPACT_SOCIAL)
    pdf.sub_title("10.2 Impact environnemental mesurable")
    pdf.kpi_table(IMPACT_ENVIRONNEMENTAL)
    pdf.sub_title("10.3 Alignement ODD")
    for odd, name, desc in ODD:
        pdf.bullet(f"{odd} - {name} : {desc}")

    # 11. Strategie
    pdf.add_page()
    pdf.section_title("11", "Strategie d'innovation et de croissance")
    pdf.sub_title("11.1 Innovation technologique")
    for i in INNOVATION:
        pdf.bullet(i)
    pdf.ln(2)
    pdf.sub_title("11.2 Strategie de croissance")
    for phase, desc in CROISSANCE:
        pdf.bullet(f"{phase} : {desc}")
    pdf.ln(2)
    pdf.sub_title("11.3 Partenariats strategiques")
    for cat, desc in PARTENARIATS:
        pdf.bullet(f"{cat} : {desc}")

    # 12. Objectifs pilote
    pdf.section_title("12", "Objectifs du pilote POESAM")
    pdf.kpi_table(OBJECTIFS_PILOTE)

    # 13. Conclusion
    pdf.section_title("13", "Conclusion")
    pdf.body_text(
        "FreshCoop represente une opportunite unique de transformer la chaine de valeur "
        "agricole senegalaise en combinant technologie, inclusion financiere et durabilite "
        "environnementale."
    )
    pdf.body_text(
        "Avec un prototype fonctionnel (web + mobile + USSD), une equipe pluridisciplinaire "
        "et des cooperatives pilotes deja identifiees, FreshCoop est pret a passer a l'echelle."
    )
    pdf.ln(4)
    pdf.set_fill_color(*GREEN)
    pdf.rect(10, pdf.get_y(), 190, 30, "F")
    pdf.set_y(pdf.get_y() + 6)
    pdf.set_text_color(*WHITE)
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "FreshCoop - Connecter. Conserver. Financer.", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 8, "contact@freshcoop.sn  |  Dakar, Senegal  |  POESAM 2026", align="C")

    out = os.path.join(DIR, "FreshCoop_Dossier_Candidature.pdf")
    pdf.output(out)
    print(f"PDF : {out}")
    return out


# ============================================================
# WORD (DOCX) GENERATION
# ============================================================

def set_cell_shading(cell, color_hex):
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)

def add_styled_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = h
        set_cell_shading(cell, "228B22")
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(9)
    for r_idx, row_data in enumerate(rows):
        row = table.rows[r_idx + 1]
        is_bold = any(k in str(row_data[0]) for k in ["TOTAL", "RESULTAT", "Tresorerie fin"])
        for c_idx, val in enumerate(row_data):
            cell = row.cells[c_idx]
            cell.text = str(val)
            p = cell.paragraphs[0]
            if c_idx > 0:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.font.size = Pt(9)
                if is_bold:
                    run.font.bold = True
            if r_idx % 2 == 1:
                set_cell_shading(cell, "E6F5E6")
    if col_widths:
        for i, w in enumerate(col_widths):
            for row in table.rows:
                row.cells[i].width = Cm(w)
    return table


def build_docx():
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = RGBColor(30, 30, 30)

    for lvl in range(1, 4):
        hs = doc.styles[f'Heading {lvl}']
        hs.font.color.rgb = RGBColor(34, 139, 34)

    # === COVER ===
    for _ in range(4):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("FRESHCOOP")
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(34, 139, 34)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("La plateforme cooperative qui connecte les producteurs aux marches")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(100, 100, 100)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Micro-hubs solaires  |  Intelligence marche  |  Preuve economique")
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(34, 139, 34)

    doc.add_paragraph()
    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Dossier de Candidature")
    run.font.size = Pt(20)
    run.font.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Programme POESAM 2026")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(100, 100, 100)

    doc.add_paragraph()
    info_items = [
        ("Porteur du projet :", "FreshCoop SAS"),
        ("Localisation :", "Dakar, Senegal"),
        ("Secteur :", "AgriTech / FinTech / Impact Social"),
        ("Date :", "Mai 2026"),
        ("Contact :", "contact@freshcoop.sn"),
    ]
    for label, val in info_items:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(label + "  ")
        run.font.color.rgb = RGBColor(100, 100, 100)
        run.font.size = Pt(11)
        run = p.add_run(val)
        run.font.bold = True
        run.font.size = Pt(11)

    doc.add_page_break()

    # === TABLE DES MATIERES ===
    doc.add_heading("Table des matieres", level=1)
    for num, title in TOC:
        p = doc.add_paragraph()
        run = p.add_run(f"{num}. ")
        run.font.bold = True
        p.add_run(title)

    doc.add_page_break()

    # === 1. Resume executif ===
    doc.add_heading("1. Resume executif", level=1)
    doc.add_paragraph(RESUME_EXECUTIF)
    doc.add_paragraph("Notre solution repose sur trois piliers complementaires :")
    for p_text in RESUME_PILIERS:
        doc.add_paragraph(p_text, style='List Bullet')
    doc.add_paragraph(RESUME_SUITE)
    doc.add_paragraph(RESUME_CANDIDATURE)

    # === 2. Probleme identifie ===
    doc.add_heading("2. Probleme identifie", level=1)
    doc.add_paragraph(PROBLEME_INTRO)
    for i, (t, d) in enumerate(PROBLEMES, 1):
        doc.add_heading(f"2.{i} {t}", level=2)
        doc.add_paragraph(d)

    doc.add_page_break()

    # === 3. Solution proposee ===
    doc.add_heading("3. Solution proposee : FreshCoop", level=1)
    doc.add_paragraph(SOLUTION_INTRO)
    for i, (t, d) in enumerate(SOLUTIONS, 1):
        doc.add_heading(f"3.{i} {t}", level=2)
        doc.add_paragraph(d)

    doc.add_page_break()

    # === 4. Analyse du marche ===
    doc.add_heading("4. Analyse du marche et des concurrents", level=1)
    doc.add_heading("4.1 Taille du marche", level=2)
    doc.add_paragraph(MARCHE_TAILLE)
    doc.add_heading("4.2 Paysage concurrentiel", level=2)
    for name, desc in CONCURRENTS:
        doc.add_paragraph(f"{name} : {desc}", style='List Bullet')
    doc.add_paragraph(
        "FreshCoop est le seul acteur qui integre stockage froid solaire, intelligence marche, "
        "tracabilite complete et preuve economique portable dans une seule plateforme."
    )
    doc.add_heading("4.3 Avantage concurrentiel", level=2)
    for a in AVANTAGES:
        doc.add_paragraph(a, style='List Bullet')

    doc.add_page_break()

    # === 5. Produit et fonctionnalites ===
    doc.add_heading("5. Produit et fonctionnalites", level=1)
    doc.add_heading("5.1 Plateforme web (React 19 + Vite + Node.js)", level=2)
    doc.add_paragraph("Application web complete avec tableau de bord adapte a chaque role.")
    doc.add_heading("5.2 Application mobile (React Native / Expo)", level=2)
    doc.add_paragraph("Application native iOS et Android avec QR code, notifications push et GPS.")
    doc.add_heading("5.3 Fonctionnalites par role", level=2)
    for role, desc in ROLES:
        p = doc.add_paragraph()
        run = p.add_run(role + " : ")
        run.font.bold = True
        run.font.color.rgb = RGBColor(34, 139, 34)
        p.add_run(desc)
    doc.add_heading("5.4 Modules transversaux", level=2)
    for mod, desc in MODULES:
        p = doc.add_paragraph()
        run = p.add_run(mod + " : ")
        run.font.bold = True
        run.font.color.rgb = RGBColor(34, 139, 34)
        p.add_run(desc)

    doc.add_page_break()

    # === 6. Equipe ===
    doc.add_heading("6. Equipe et competences cles", level=1)
    doc.add_paragraph("L'equipe FreshCoop reunit des competences complementaires :")
    for role, desc in EQUIPE:
        p = doc.add_paragraph()
        run = p.add_run(role)
        run.font.bold = True
        run.font.color.rgb = RGBColor(34, 139, 34)
        p.add_run(f"\n{desc}")

    # === 7. Modele economique ===
    doc.add_heading("7. Modele economique", level=1)
    doc.add_paragraph("FreshCoop adopte un modele de revenus diversifie :")
    for r in MODELE_REVENUS:
        doc.add_paragraph(r, style='List Bullet')
    doc.add_paragraph(
        "Point cle : FreshCoop n'est PAS un portefeuille electronique. Les paiements sont "
        "orchestres par des partenaires agrees."
    )

    doc.add_page_break()

    # === 8. Previsions financieres ===
    doc.add_heading("8. Previsions financieres sur 3 ans", level=1)
    doc.add_heading("8.1 Hypotheses cles", level=2)
    doc.add_paragraph("Annee 1 (Pilote) : 100 producteurs, 3 cooperatives, 4 regions", style='List Bullet')
    doc.add_paragraph("Annee 2 (Croissance) : 500 producteurs, 12 cooperatives, 6 regions", style='List Bullet')
    doc.add_paragraph("Annee 3 (Echelle) : 2 000 producteurs, 40 cooperatives, 10 regions", style='List Bullet')

    doc.add_heading("8.2 Compte de resultat previsionnel (M FCFA)", level=2)
    add_styled_table(doc, ["Poste", "Annee 1", "Annee 2", "Annee 3"], PREVISIONS)

    doc.add_paragraph()
    doc.add_heading("8.3 Tresorerie previsionnelle (M FCFA)", level=2)
    add_styled_table(doc, ["Poste", "Annee 1", "Annee 2", "Annee 3"], TRESORERIE)
    doc.add_paragraph("Le break-even est atteint en Annee 3 avec une marge nette de 38%.")

    doc.add_page_break()

    # === 9. Utilisation du financement ===
    doc.add_heading("9. Utilisation prevue du financement", level=1)
    doc.add_paragraph("Le financement POESAM (50 millions FCFA) sera alloue comme suit :")
    add_styled_table(doc, ["Poste", "Montant", "Part"], FINANCEMENT_REPARTITION)

    # === 10. Impact ===
    doc.add_heading("10. Impact social et environnemental", level=1)
    doc.add_heading("10.1 Impact social mesurable", level=2)
    add_styled_table(doc, ["Indicateur", "Objectif", "Horizon"], IMPACT_SOCIAL)
    doc.add_paragraph()
    doc.add_heading("10.2 Impact environnemental mesurable", level=2)
    add_styled_table(doc, ["Indicateur", "Objectif", "Horizon"], IMPACT_ENVIRONNEMENTAL)
    doc.add_paragraph()
    doc.add_heading("10.3 Alignement ODD", level=2)
    for odd, name, desc in ODD:
        doc.add_paragraph(f"{odd} - {name} : {desc}", style='List Bullet')

    doc.add_page_break()

    # === 11. Strategie ===
    doc.add_heading("11. Strategie d'innovation et de croissance", level=1)
    doc.add_heading("11.1 Innovation technologique", level=2)
    for i in INNOVATION:
        doc.add_paragraph(i, style='List Bullet')
    doc.add_heading("11.2 Strategie de croissance", level=2)
    for phase, desc in CROISSANCE:
        doc.add_paragraph(f"{phase} : {desc}", style='List Bullet')
    doc.add_heading("11.3 Partenariats strategiques", level=2)
    for cat, desc in PARTENARIATS:
        doc.add_paragraph(f"{cat} : {desc}", style='List Bullet')

    # === 12. Objectifs pilote ===
    doc.add_heading("12. Objectifs du pilote POESAM", level=1)
    add_styled_table(doc, ["Indicateur", "Objectif", "Horizon"], OBJECTIFS_PILOTE)

    # === 13. Conclusion ===
    doc.add_heading("13. Conclusion", level=1)
    doc.add_paragraph(
        "FreshCoop represente une opportunite unique de transformer la chaine de valeur "
        "agricole senegalaise en combinant technologie, inclusion financiere et durabilite "
        "environnementale."
    )
    doc.add_paragraph(
        "Avec un prototype fonctionnel (web + mobile + USSD), une equipe pluridisciplinaire "
        "et des cooperatives pilotes deja identifiees, FreshCoop est pret a passer a l'echelle "
        "avec le soutien de POESAM."
    )
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("FreshCoop - Connecter. Conserver. Financer.")
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(34, 139, 34)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("contact@freshcoop.sn  |  Dakar, Senegal  |  POESAM 2026")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(100, 100, 100)

    out = os.path.join(DIR, "FreshCoop_Dossier_Candidature.docx")
    doc.save(out)
    print(f"DOCX : {out}")
    return out


# ============================================================
# POWERPOINT (PPTX) GENERATION
# ============================================================

def add_green_bar(slide, top=0, height=PptxInches(0.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        PptxInches(0), top,
        PptxInches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptxRGBColor(34, 139, 34)
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = PptxRGBColor(*color)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=(30, 30, 30)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = PptxPt(font_size)
        p.font.name = "Calibri"
        p.font.color.rgb = PptxRGBColor(*color)
        p.space_after = PptxPt(6)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>')
        pPr.append(buChar)
    return txBox

def add_pptx_table(slide, left, top, width, height, headers, rows):
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = PptxPt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PptxRGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptxRGBColor(34, 139, 34)
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(9)
                if c_idx > 0:
                    paragraph.alignment = PP_ALIGN.CENTER
                is_bold = any(k in str(row_data[0]) for k in ["TOTAL", "RESULTAT", "Tresorerie fin"])
                if is_bold:
                    paragraph.font.bold = True
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGBColor(230, 245, 230)
    return table_shape


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(blank)
    add_green_bar(slide, PptxInches(1.5), PptxInches(2.5))
    add_text_box(slide, PptxInches(0), PptxInches(1.8), PptxInches(13.333), PptxInches(1),
                 "FRESHCOOP", 48, True, (255, 255, 255), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(2.5), PptxInches(13.333), PptxInches(0.6),
                 "La plateforme cooperative qui connecte les producteurs aux marches",
                 16, False, (255, 255, 255), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(3.0), PptxInches(13.333), PptxInches(0.5),
                 "Micro-hubs solaires  |  Intelligence marche  |  Preuve economique",
                 12, True, (255, 255, 255), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(4.5), PptxInches(13.333), PptxInches(0.6),
                 "Dossier de Candidature - Programme POESAM 2026",
                 22, True, (30, 30, 30), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(5.3), PptxInches(13.333), PptxInches(0.5),
                 "FreshCoop SAS  |  Dakar, Senegal  |  AgriTech / FinTech / Impact Social  |  Mai 2026",
                 12, False, (100, 100, 100), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(5.8), PptxInches(13.333), PptxInches(0.4),
                 "contact@freshcoop.sn",
                 11, False, (100, 100, 100), PP_ALIGN.CENTER)

    # ── Slide 2: Sommaire ──
    slide = prs.slides.add_slide(blank)
    add_green_bar(slide, PptxInches(0), PptxInches(0.6))
    add_text_box(slide, PptxInches(0.3), PptxInches(0.6), PptxInches(6), PptxInches(0.5),
                 "SOMMAIRE", 24, True, (255, 255, 255), PP_ALIGN.LEFT)
    items_left = [f"{n}. {t}" for n, t in TOC[:7]]
    items_right = [f"{n}. {t}" for n, t in TOC[7:]]
    add_bullet_list(slide, PptxInches(0.5), PptxInches(1.5), PptxInches(6), PptxInches(5), items_left, 14)
    add_bullet_list(slide, PptxInches(6.8), PptxInches(1.5), PptxInches(6), PptxInches(5), items_right, 14)

    # Helper for section slides
    def section_slide(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        add_green_bar(s, PptxInches(0), PptxInches(0.6))
        add_text_box(s, PptxInches(0.3), PptxInches(0.6), PptxInches(12), PptxInches(0.5),
                     title.upper(), 22, True, (255, 255, 255), PP_ALIGN.LEFT)
        if subtitle:
            add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(12), PptxInches(0.6),
                         subtitle, 14, False, (100, 100, 100))
        # Footer
        add_text_box(s, PptxInches(0), PptxInches(7.0), PptxInches(13.333), PptxInches(0.4),
                     "FreshCoop  |  Dossier POESAM 2026  |  Confidentiel",
                     8, False, (150, 150, 150), PP_ALIGN.CENTER)
        return s

    # ── Slide 3: Resume executif ──
    s = section_slide("1. Resume executif")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(12), PptxInches(1.2),
                 RESUME_EXECUTIF, 13, False, (30, 30, 30))
    add_text_box(s, PptxInches(0.5), PptxInches(2.6), PptxInches(12), PptxInches(0.4),
                 "Trois piliers complementaires :", 14, True, (34, 139, 34))
    add_bullet_list(s, PptxInches(0.5), PptxInches(3.1), PptxInches(12), PptxInches(1.5),
                    RESUME_PILIERS, 13)
    add_text_box(s, PptxInches(0.5), PptxInches(4.8), PptxInches(12), PptxInches(1.5),
                 RESUME_SUITE, 12, False, (60, 60, 60))

    # ── Slide 4: Probleme ──
    s = section_slide("2. Probleme identifie", PROBLEME_INTRO)
    y = 2.2
    for i, (t, d) in enumerate(PROBLEMES):
        # Create a colored box for each problem
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   PptxInches(0.5), PptxInches(y),
                                   PptxInches(12.3), PptxInches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(245, 255, 245)
        shape.line.color.rgb = PptxRGBColor(34, 139, 34)
        shape.line.width = PptxPt(1)
        add_text_box(s, PptxInches(0.8), PptxInches(y + 0.1), PptxInches(11.5), PptxInches(0.3),
                     t, 14, True, (34, 139, 34))
        add_text_box(s, PptxInches(0.8), PptxInches(y + 0.5), PptxInches(11.5), PptxInches(0.7),
                     d, 11, False, (60, 60, 60))
        y += 1.5

    # ── Slide 5: Solution ──
    s = section_slide("3. Solution proposee : FreshCoop")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(12), PptxInches(0.5),
                 SOLUTION_INTRO, 13, False, (60, 60, 60))
    col_w = 3.0
    for i, (t, d) in enumerate(SOLUTIONS):
        x = 0.4 + i * (col_w + 0.15)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   PptxInches(x), PptxInches(2.2),
                                   PptxInches(col_w), PptxInches(4.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(245, 255, 245)
        shape.line.color.rgb = PptxRGBColor(34, 139, 34)
        shape.line.width = PptxPt(1)
        add_text_box(s, PptxInches(x + 0.15), PptxInches(2.4), PptxInches(col_w - 0.3), PptxInches(0.5),
                     t, 13, True, (34, 139, 34), PP_ALIGN.CENTER)
        add_text_box(s, PptxInches(x + 0.15), PptxInches(3.0), PptxInches(col_w - 0.3), PptxInches(3.2),
                     d, 10, False, (60, 60, 60))

    # ── Slide 6: Marche ──
    s = section_slide("4. Analyse du marche")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(5.5), PptxInches(0.4),
                 "Taille du marche", 16, True, (34, 139, 34))
    add_text_box(s, PptxInches(0.5), PptxInches(1.8), PptxInches(5.5), PptxInches(1.5),
                 MARCHE_TAILLE, 11, False, (60, 60, 60))

    # Key figures
    figures = [
        ("3 000 Mds", "FCFA/an\nMarche agricole"),
        ("800 Mds", "FCFA\nFruits & legumes"),
        ("40%", "Pertes\npost-recolte"),
        ("300 Mds", "FCFA\nValeur detruite"),
    ]
    for i, (val, label) in enumerate(figures):
        x = 0.5 + i * 3.1
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   PptxInches(x), PptxInches(3.5),
                                   PptxInches(2.8), PptxInches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(34, 139, 34)
        shape.line.fill.background()
        add_text_box(s, PptxInches(x), PptxInches(3.55), PptxInches(2.8), PptxInches(0.5),
                     val, 22, True, (255, 255, 255), PP_ALIGN.CENTER)
        add_text_box(s, PptxInches(x), PptxInches(4.05), PptxInches(2.8), PptxInches(0.6),
                     label, 10, False, (255, 255, 255), PP_ALIGN.CENTER)

    add_text_box(s, PptxInches(0.5), PptxInches(5.0), PptxInches(12), PptxInches(0.4),
                 "Avantages concurrentiels", 14, True, (34, 139, 34))
    add_bullet_list(s, PptxInches(0.5), PptxInches(5.5), PptxInches(12), PptxInches(1.5),
                    AVANTAGES[:3], 11)

    # ── Slide 7: Produit ──
    s = section_slide("5. Produit et fonctionnalites")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(12), PptxInches(0.5),
                 "Web (React 19 + Vite)  |  Mobile (React Native / Expo)  |  USSD (*384*FRES#)",
                 14, True, (34, 139, 34))
    items = [f"{r}: {d}" for r, d in ROLES]
    add_bullet_list(s, PptxInches(0.5), PptxInches(2.0), PptxInches(6), PptxInches(5), items[:4], 11)
    add_bullet_list(s, PptxInches(6.8), PptxInches(2.0), PptxInches(6), PptxInches(5), items[4:], 11)

    # ── Slide 8: Equipe ──
    s = section_slide("6. Equipe et competences cles")
    col_w = 4.0
    for i, (role, desc) in enumerate(EQUIPE):
        col = i % 3
        row = i // 3
        x = 0.4 + col * (col_w + 0.15)
        y_start = 1.5 + row * 2.6
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   PptxInches(x), PptxInches(y_start),
                                   PptxInches(col_w), PptxInches(2.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(245, 255, 245)
        shape.line.color.rgb = PptxRGBColor(34, 139, 34)
        shape.line.width = PptxPt(1)
        add_text_box(s, PptxInches(x + 0.2), PptxInches(y_start + 0.2),
                     PptxInches(col_w - 0.4), PptxInches(0.5),
                     role, 12, True, (34, 139, 34))
        add_text_box(s, PptxInches(x + 0.2), PptxInches(y_start + 0.8),
                     PptxInches(col_w - 0.4), PptxInches(1.3),
                     desc, 10, False, (60, 60, 60))

    # ── Slide 9: Modele economique ──
    s = section_slide("7. Modele economique")
    add_bullet_list(s, PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(3),
                    MODELE_REVENUS, 13)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               PptxInches(0.5), PptxInches(5.0),
                               PptxInches(12.3), PptxInches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptxRGBColor(255, 250, 230)
    shape.line.color.rgb = PptxRGBColor(200, 150, 0)
    add_text_box(s, PptxInches(0.8), PptxInches(5.1), PptxInches(11.8), PptxInches(1),
                 "Point cle : FreshCoop n'est PAS un portefeuille electronique. "
                 "Les paiements passent par Orange Money, Wave, PayDunya. "
                 "Aucun fonds detenu = zero contrainte reglementaire EME.",
                 12, False, (80, 60, 0))

    # ── Slide 10: Previsions financieres ──
    s = section_slide("8. Previsions financieres sur 3 ans")
    add_pptx_table(s,
                   PptxInches(0.5), PptxInches(1.5), PptxInches(12.3), PptxInches(4.5),
                   ["Poste", "Annee 1 (M FCFA)", "Annee 2 (M FCFA)", "Annee 3 (M FCFA)"],
                   PREVISIONS)

    # ── Slide 11: Tresorerie + financement ──
    s = section_slide("9. Tresorerie et utilisation du financement")
    add_pptx_table(s,
                   PptxInches(0.3), PptxInches(1.4), PptxInches(6.2), PptxInches(2.2),
                   ["Poste", "An 1", "An 2", "An 3"],
                   TRESORERIE)
    add_text_box(s, PptxInches(7), PptxInches(1.4), PptxInches(6), PptxInches(0.4),
                 "Repartition du financement (50 M FCFA)", 14, True, (34, 139, 34))
    add_pptx_table(s,
                   PptxInches(7), PptxInches(2.0), PptxInches(6), PptxInches(2.8),
                   ["Poste", "Montant", "Part"],
                   FINANCEMENT_REPARTITION)

    # ── Slide 12: Impact ──
    s = section_slide("10. Impact social et environnemental")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(6), PptxInches(0.4),
                 "Impact social", 16, True, (34, 139, 34))
    add_pptx_table(s,
                   PptxInches(0.3), PptxInches(1.8), PptxInches(6.2), PptxInches(3),
                   ["Indicateur", "Objectif", "Horizon"],
                   IMPACT_SOCIAL)
    add_text_box(s, PptxInches(7), PptxInches(1.3), PptxInches(6), PptxInches(0.4),
                 "Impact environnemental", 16, True, (34, 139, 34))
    add_pptx_table(s,
                   PptxInches(7), PptxInches(1.8), PptxInches(6), PptxInches(2.5),
                   ["Indicateur", "Objectif", "Horizon"],
                   IMPACT_ENVIRONNEMENTAL)
    # ODD
    add_text_box(s, PptxInches(0.5), PptxInches(5.2), PptxInches(12), PptxInches(0.4),
                 "Alignement avec les Objectifs de Developpement Durable (ODD)", 13, True, (34, 139, 34))
    odd_items = [f"{o} - {n} : {d}" for o, n, d in ODD]
    add_bullet_list(s, PptxInches(0.5), PptxInches(5.7), PptxInches(12), PptxInches(1.5), odd_items, 10)

    # ── Slide 13: Strategie ──
    s = section_slide("11. Strategie d'innovation et de croissance")
    add_text_box(s, PptxInches(0.5), PptxInches(1.3), PptxInches(6), PptxInches(0.4),
                 "Innovation technologique", 14, True, (34, 139, 34))
    add_bullet_list(s, PptxInches(0.5), PptxInches(1.8), PptxInches(5.5), PptxInches(2), INNOVATION, 11)

    add_text_box(s, PptxInches(7), PptxInches(1.3), PptxInches(6), PptxInches(0.4),
                 "Strategie de croissance", 14, True, (34, 139, 34))
    growth_items = [f"{p} : {d}" for p, d in CROISSANCE]
    add_bullet_list(s, PptxInches(7), PptxInches(1.8), PptxInches(5.5), PptxInches(2), growth_items, 11)

    add_text_box(s, PptxInches(0.5), PptxInches(4.2), PptxInches(12), PptxInches(0.4),
                 "Partenariats strategiques", 14, True, (34, 139, 34))
    partner_items = [f"{c} : {d}" for c, d in PARTENARIATS]
    add_bullet_list(s, PptxInches(0.5), PptxInches(4.7), PptxInches(12), PptxInches(2), partner_items, 11)

    # ── Slide 14: Objectifs pilote ──
    s = section_slide("12. Objectifs du pilote POESAM")
    add_pptx_table(s,
                   PptxInches(1.5), PptxInches(1.8), PptxInches(10), PptxInches(3.5),
                   ["Indicateur", "Objectif", "Horizon"],
                   OBJECTIFS_PILOTE)

    # ── Slide 15: Conclusion ──
    slide = prs.slides.add_slide(blank)
    add_green_bar(slide, PptxInches(2.5), PptxInches(3))
    add_text_box(slide, PptxInches(0), PptxInches(1.0), PptxInches(13.333), PptxInches(1),
                 "CONCLUSION", 32, True, (34, 139, 34), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(1), PptxInches(1.8), PptxInches(11.333), PptxInches(0.8),
                 "FreshCoop represente une opportunite unique de transformer la chaine de valeur "
                 "agricole senegalaise en combinant technologie, inclusion financiere et durabilite.",
                 14, False, (60, 60, 60), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(2.7), PptxInches(13.333), PptxInches(0.7),
                 "FreshCoop - Connecter. Conserver. Financer.",
                 28, True, (255, 255, 255), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(3.4), PptxInches(13.333), PptxInches(0.5),
                 "Micro-hubs solaires  |  Intelligence marche  |  Preuve economique",
                 14, False, (255, 255, 255), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(5.0), PptxInches(13.333), PptxInches(0.6),
                 "Pret a passer a l'echelle avec le soutien de POESAM.",
                 18, True, (30, 30, 30), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(6.0), PptxInches(13.333), PptxInches(0.3),
                 "contact@freshcoop.sn  |  Dakar, Senegal  |  POESAM 2026",
                 11, False, (100, 100, 100), PP_ALIGN.CENTER)
    add_text_box(slide, PptxInches(0), PptxInches(6.5), PptxInches(13.333), PptxInches(0.3),
                 "FreshCoop  |  Dossier POESAM 2026  |  Confidentiel",
                 8, False, (150, 150, 150), PP_ALIGN.CENTER)

    out = os.path.join(DIR, "FreshCoop_Dossier_Candidature.pptx")
    prs.save(out)
    print(f"PPTX : {out}")
    return out


# ============================================================
# MAIN
# ============================================================

if __name__ == "__main__":
    print("=" * 60)
    print("  FreshCoop - Generation du dossier de candidature POESAM")
    print("=" * 60)
    print()
    pdf_path = build_pdf()
    docx_path = build_docx()
    pptx_path = build_pptx()
    print()
    print("=" * 60)
    print("  GENERATION TERMINEE")
    print(f"  PDF  : {pdf_path}")
    print(f"  DOCX : {docx_path}")
    print(f"  PPTX : {pptx_path}")
    print("=" * 60)
